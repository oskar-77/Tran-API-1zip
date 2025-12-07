"""PowerPoint document extractor using python-pptx."""

from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Inches
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.extractors.base_extractor import BaseExtractor
from src.schemas.document import (
    Document, Page, Metadata, TextDirection,
    HeadingBlock, ParagraphBlock, ImageBlock, TableBlock,
    TableRow, TableCell, Link
)
from src.utils.text_utils import detect_text_direction, clean_text
from src.utils.image_utils import image_to_base64, generate_image_id


class PptxExtractor(BaseExtractor):
    """Extractor for PowerPoint documents using python-pptx."""
    
    SUPPORTED_EXTENSIONS = ['.pptx']
    
    def extract(self) -> Document:
        """Extract content from PowerPoint document."""
        prs = Presentation(str(self.file_path))
        
        metadata = self._extract_metadata(prs)
        pages = []
        
        for slide_index, slide in enumerate(prs.slides):
            page = self._extract_slide(slide, slide_index + 1)
            pages.append(page)
        
        all_text = " ".join([
            block.text for page in pages for block in page.blocks 
            if hasattr(block, 'text')
        ])
        direction = TextDirection(detect_text_direction(all_text))
        
        metadata.page_count = len(pages)
        
        return Document(
            title=metadata.title or self.file_path.stem,
            metadata=metadata,
            pages=pages,
            direction=direction
        )
    
    def _extract_metadata(self, prs: Presentation) -> Metadata:
        """Extract metadata from PowerPoint."""
        base_metadata = self._create_base_metadata()
        core_props = prs.core_properties
        
        return Metadata(
            title=core_props.title or base_metadata.source_filename,
            author=core_props.author,
            created_date=core_props.created,
            modified_date=core_props.modified or base_metadata.modified_date,
            subject=core_props.subject,
            keywords=self._parse_keywords(core_props.keywords),
            source_format='pptx',
            source_filename=base_metadata.source_filename
        )
    
    def _parse_keywords(self, keywords_str: Optional[str]) -> list[str]:
        """Parse keywords string into list."""
        if not keywords_str:
            return []
        
        keywords = [k.strip() for k in keywords_str.split(',')]
        return [k for k in keywords if k]
    
    def _extract_slide(self, slide, slide_num: int) -> Page:
        """Extract content from a slide."""
        blocks = []
        slide_title = None
        
        for shape in slide.shapes:
            if shape.is_placeholder:
                if hasattr(shape, "text") and shape.text:
                    placeholder_type = shape.placeholder_format.type
                    
                    if placeholder_type == 1:
                        slide_title = clean_text(shape.text)
                        blocks.append(HeadingBlock(
                            level=1,
                            text=slide_title,
                            direction=TextDirection(detect_text_direction(slide_title)),
                            position_hint=f"slide_{slide_num}_title"
                        ))
                    else:
                        text_blocks = self._extract_text_frame(shape.text_frame, slide_num)
                        blocks.extend(text_blocks)
            
            elif shape.has_text_frame:
                text_blocks = self._extract_text_frame(shape.text_frame, slide_num)
                blocks.extend(text_blocks)
            
            if shape.has_table:
                table_block = self._extract_table(shape.table, slide_num)
                if table_block:
                    blocks.append(table_block)
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_block = self._extract_image(shape, slide_num)
                if image_block:
                    blocks.append(image_block)
        
        all_text = " ".join([b.text for b in blocks if hasattr(b, 'text')])
        direction = TextDirection(detect_text_direction(all_text))
        
        return Page(
            page_number=slide_num,
            title=slide_title,
            blocks=blocks,
            direction=direction
        )
    
    def _extract_text_frame(self, text_frame, slide_num: int) -> list:
        """Extract text blocks from text frame."""
        blocks = []
        
        for para in text_frame.paragraphs:
            text = clean_text(para.text)
            if not text:
                continue
            
            direction = TextDirection(detect_text_direction(text))
            
            level = para.level or 0
            if level == 0 and len(text) < 100:
                blocks.append(HeadingBlock(
                    level=min(level + 2, 6),
                    text=text,
                    direction=direction,
                    position_hint=f"slide_{slide_num}"
                ))
            else:
                is_bold = any(run.font.bold for run in para.runs if run.font.bold)
                is_italic = any(run.font.italic for run in para.runs if run.font.italic)
                
                blocks.append(ParagraphBlock(
                    text=text,
                    direction=direction,
                    is_bold=is_bold,
                    is_italic=is_italic,
                    position_hint=f"slide_{slide_num}"
                ))
        
        return blocks
    
    def _extract_table(self, table, slide_num: int) -> Optional[TableBlock]:
        """Extract table from slide."""
        rows = []
        
        for row_index, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = clean_text(cell.text)
                cells.append(TableCell(
                    content=cell_text,
                    is_header=(row_index == 0)
                ))
            
            rows.append(TableRow(
                cells=cells,
                is_header_row=(row_index == 0)
            ))
        
        if not rows:
            return None
        
        return TableBlock(
            rows=rows,
            has_header=True,
            row_count=len(rows),
            column_count=len(rows[0].cells) if rows else 0,
            position_hint=f"slide_{slide_num}"
        )
    
    def _extract_image(self, shape, slide_num: int) -> Optional[ImageBlock]:
        """Extract image from shape."""
        try:
            image = shape.image
            image_data = image.blob
            image_format = image.ext
            
            image_id = generate_image_id()
            base64_data = image_to_base64(image_data, image_format)
            
            return ImageBlock(
                image_id=image_id,
                image_data=base64_data,
                image_format=image_format,
                position_hint=f"slide_{slide_num}"
            )
        except Exception:
            return None
